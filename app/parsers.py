from pathlib import Path
from typing import List, Dict
import csv
def parse_txt_md(path: str) -> List[Dict]:
    content = Path(path).read_text(encoding='utf-8')
    return [{"text": content}]

# Lightweight CSV parser
def parse_csv(path: str) -> List[Dict]:
    rows = []
    with open(path, newline='', encoding='utf-8') as f:
        reader = csv.DictReader(f)
        for i, r in enumerate(reader):
            rows.append({"row_index": i, "text": str(r)})
    return rows

# Placeholder parsers for PDF/PPTX/DOCX — they require external libs
def parse_pdf(path: str) -> List[Dict]:
    # Requires pdfplumber; simple fallback to txt if not installed
    try:
        import pdfplumber
        out = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                out.append({"page": i+1, "text": page.extract_text() or ""})
        return out
    except Exception:
        return parse_txt_md(path)

def parse_pptx(path: str) -> List[Dict]:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            txt = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt.append(shape.text)
            slides.append({"slide": i+1, "text": "\n".join(txt)})
        return slides
    except Exception:
        return parse_txt_md(path)

def parse_docx(path: str) -> List[Dict]:
    try:
        from docx import Document
        doc = Document(path)
        paras = []
        for i, p in enumerate(doc.paragraphs):
            paras.append({"para_index": i, "text": p.text})
        return paras
    except Exception:
        return parse_txt_md(path)
